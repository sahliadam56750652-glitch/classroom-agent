"""Turn fetched bytes into plain text.

Dispatch is on the mime type Drive reported, never on the filename. Teachers
post files whose titles disagree with their contents, and one tracked file is a
legacy `application/vnd.ms-powerpoint` that python-pptx cannot open at all --
branching on the extension would hand it to the wrong reader and get back a
confusing PackageNotFoundError instead of an honest "unsupported".

Nothing here calls an LLM, and nothing raises on a file it cannot read: an
unreadable format is a recorded row and the run continues.

On PDFs and OCR
---------------
The decision is made per PAGE, not per file. A document can be typed notes with
scanned exercises stapled on the end, and both halves are common in this
account. Each page is classified:

  native  enough text to be worth keeping
  scan    almost no text, but the page carries an image -- OCR would reach it
  blank   almost no text and nothing drawn on it -- a section divider

Blank pages are the reason the rule is not simply "low text means OCR it".
Sending them to OCR would multiply the work for guaranteed nothing.

The OCR step itself is deliberately NOT wired up yet. The evidence for needing
it -- 3 of 5 sampled PDFs -- was measured on a course this sync does not track,
so the real figure for the tracked library is unknown. `scan_pages` records what
OCR would have to read, `agent extract --dry-run` totals it, and that number
decides whether installing Tesseract is worth it.
"""

from __future__ import annotations

import json
import sqlite3
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import pymupdf

from ..config import Config
from ..db import store

TEXT_SUBDIR = "text"

# Page separator in the written .txt. A form feed is what a page break has meant
# in plain text for fifty years, and it survives a round trip through anything
# that handles text at all.
PAGE_BREAK = "\f"

# Below this many characters a page is not carrying its content as text. The
# threshold comes from the Phase 0 probe: the native samples measured 285 and
# 598 chars/page, the scanned ones 0, 11.3 and 19 -- an order of magnitude
# apart, so the exact value is not delicate.
MIN_CHARS_PER_PAGE = 100

PDF_MIME = "application/pdf"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
NOTEBOOK_MIMES = frozenset({"application/json", "application/octet-stream"})

# An image posted as an attachment carries no text layer at all, so it is 100%
# OCR by definition. It is recorded here as one page that is one scan page,
# which is exactly what a PDF page hiding a diagram looks like -- so files/ocr.py
# picks it up through the same scan_pages route and needs no special case.
IMAGE_MIMES = frozenset({"image/png", "image/jpeg"})

# Anything whose bytes are text already. text/x-sql and the notebooks are source
# code rather than prose; they are extracted because they are course material,
# and `method` records what they are so Phase 3 can decline to quiz on them.
PLAIN_TEXT_MIMES = frozenset(
    {
        "text/plain",
        "text/csv",
        "text/x-sql",
        "text/x-python",
        "text/markdown",
        "text/html",
    }
)

# Formats we can fetch but not read, with the reason recorded rather than a
# silent skip. Legacy .ppt and .doc are OLE compound files, not zip packages,
# and the modern readers cannot open them.
UNREADABLE = {
    "application/zip": "zip archive -- not unpacked",
    "application/x-zip-compressed": "zip archive -- not unpacked",
    "application/vnd.ms-powerpoint": "legacy .ppt -- python-pptx reads only .pptx",
    "application/msword": "legacy .doc -- python-docx reads only .docx",
    "application/seb": "Safe Exam Browser config -- no text content",
    "video/mp4": "video -- no transcription in this phase",
}

# Encodings tried in order. latin-1 cannot fail, so it terminates the chain and
# guarantees a decode rather than an exception that kills the run over one file.
ENCODINGS = ("utf-8-sig", "cp1252", "latin-1")


class ExtractError(Exception):
    """Extraction is impossible for a reason that is not about one file."""


@dataclass
class PageStats:
    native: int = 0
    scan: int = 0
    blank: int = 0

    @property
    def total(self) -> int:
        return self.native + self.scan + self.blank


@dataclass
class Extracted:
    """One file's text and how it was obtained."""

    text: str
    method: str
    pages: int | None = None
    scan_pages: int = 0
    ocr_pages: int = 0


@dataclass
class ExtractResult:
    candidates: int = 0
    extracted: int = 0
    unsupported: int = 0
    errors: int = 0
    skipped: int = 0
    chars: int = 0
    scan_pages: int = 0
    pdf_pages: int = 0
    dry_run: bool = False
    # (title, chars_per_page) for every PDF seen, so the dry run can print the
    # distribution that decides the OCR question.
    pdf_profile: list[tuple[str, float, int]] = field(default_factory=list)
    reasons: dict[str, int] = field(default_factory=dict)

    def items_seen(self) -> dict[str, int]:
        return {
            "candidates": self.candidates,
            "extracted": self.extracted,
            "unsupported": self.unsupported,
            "errors": self.errors,
            "skipped": self.skipped,
            "scan_pages": self.scan_pages,
        }


def decode(payload: bytes) -> str:
    """Text bytes to str, without ever raising.

    The 14 text/plain files in the tracked courses have unknown encodings; a
    bare utf-8 decode dies on the first cp1252 apostrophe and takes the run
    with it.
    """
    for encoding in ENCODINGS:
        try:
            return payload.decode(encoding)
        except UnicodeDecodeError:
            continue
    return payload.decode("latin-1", errors="replace")  # pragma: no cover - latin-1 cannot fail


def classify_page(page: Any) -> tuple[str, str]:
    """(kind, text) for one PDF page: 'native', 'scan' or 'blank'."""
    text = page.get_text("text") or ""
    if len(text.strip()) >= MIN_CHARS_PER_PAGE:
        return "native", text

    # Low text. Something drawn on the page means the content is there but not
    # as text; nothing drawn means the page really is empty.
    has_image = bool(page.get_images(full=True))
    has_drawing = False
    if not has_image:
        try:
            has_drawing = bool(page.get_drawings())
        except Exception:  # pragma: no cover - defensive, varies by MuPDF build
            has_drawing = False
    return ("scan" if has_image or has_drawing else "blank"), text


def extract_pdf(path: Path) -> Extracted:
    pages: list[str] = []
    stats = PageStats()
    with pymupdf.open(path) as document:
        for page in document:
            kind, text = classify_page(page)
            setattr(stats, kind, getattr(stats, kind) + 1)
            pages.append(text)
    return Extracted(
        text=PAGE_BREAK.join(pages),
        method="pymupdf",
        pages=stats.total,
        scan_pages=stats.scan,
    )


def extract_docx(path: Path) -> Extracted:
    from docx import Document  # imported lazily so a stale venv is one bad row

    document = Document(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs]
    # Tables carry real content in lab sheets and are invisible to `paragraphs`.
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    return Extracted(text="\n".join(parts), method="docx")


def extract_pptx(path: Path) -> Extracted:
    from pptx import Presentation

    slides: list[str] = []
    for slide in Presentation(str(path)).slides:
        parts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        notes = slide.notes_slide if slide.has_notes_slide else None
        if notes is not None and notes.notes_text_frame.text.strip():
            parts.append(f"[notes] {notes.notes_text_frame.text}")
        slides.append("\n".join(parts))
    return Extracted(text=PAGE_BREAK.join(slides), method="pptx", pages=len(slides))


def extract_notebook(payload: bytes) -> Extracted:
    """Markdown and code cells from a .ipynb. Outputs are noise and are dropped."""
    document = json.loads(decode(payload))
    if not isinstance(document, dict) or "cells" not in document:
        raise ValueError("not a Jupyter notebook")
    parts: list[str] = []
    for cell in document.get("cells") or []:
        source = cell.get("source")
        body = "".join(source) if isinstance(source, list) else str(source or "")
        if body.strip():
            parts.append(body)
    return Extracted(text="\n\n".join(parts), method="ipynb", pages=len(parts))


def extract_file(path: Path, mime_type: str | None) -> Extracted | None:
    """Text for one local file, or None when the format is not readable."""
    if mime_type in UNREADABLE:
        return None
    if mime_type == PDF_MIME:
        return extract_pdf(path)
    if mime_type == DOCX_MIME:
        return extract_docx(path)
    if mime_type == PPTX_MIME:
        return extract_pptx(path)
    if mime_type in PLAIN_TEXT_MIMES:
        return Extracted(text=decode(path.read_bytes()), method="text")
    if mime_type in IMAGE_MIMES:
        # No text to extract, but this is not "unsupported" -- calling it that
        # would drop a photographed board out of the library permanently. It is
        # one page, entirely unread, waiting for `agent ocr`.
        return Extracted(text="", method="image", pages=1, scan_pages=1)
    if mime_type in NOTEBOOK_MIMES:
        # Two tracked files are notebooks that Drive typed as json/octet-stream.
        # Try it, and fall through to unsupported if it is something else.
        try:
            return extract_notebook(path.read_bytes())
        except (ValueError, UnicodeDecodeError):
            return None
    return None


def _reason(result: ExtractResult, text: str) -> None:
    result.reasons[text] = result.reasons.get(text, 0) + 1


def extract(
    config: Config,
    db: sqlite3.Connection,
    *,
    dry_run: bool = False,
    force: bool = False,
    now: str | None = None,
) -> ExtractResult:
    """Read every fetched file that has no current text.

    dry_run reads and measures but writes neither the .txt files nor the rows.
    That makes it the measurement pass: it prints the chars-per-page
    distribution across the library, which is what decides the OCR question.
    """
    rows = db.execute(
        "SELECT * FROM extractions WHERE status IN ('fetched', 'ok') "
        "AND local_path IS NOT NULL ORDER BY drive_id"
    ).fetchall()

    result = ExtractResult(dry_run=dry_run, candidates=len(rows))
    stamp = now or store._utc_now_iso()
    text_dir = config.library_dir / TEXT_SUBDIR

    titles = {
        row["drive_id"]: row["title"]
        for row in db.execute(
            "SELECT DISTINCT drive_id, title FROM materials WHERE drive_id IS NOT NULL"
        ).fetchall()
    }

    for row in rows:
        drive_id = row["drive_id"]
        if row["status"] == "ok" and row["text_path"] and not force:
            result.skipped += 1
            continue

        source = config.library_dir / row["local_path"]
        if not source.exists():
            if not dry_run:
                store.upsert_extraction(
                    db,
                    drive_id,
                    status="error",
                    error=f"local file missing: {row['local_path']}",
                    extracted_at=stamp,
                )
            result.errors += 1
            _reason(result, "local file missing")
            continue

        try:
            extracted = extract_file(source, row["mime_type"])
        except Exception as err:
            # One malformed PDF must not end the run. The row records why.
            if not dry_run:
                store.upsert_extraction(
                    db, drive_id, status="error", error=repr(err), extracted_at=stamp
                )
            result.errors += 1
            _reason(result, type(err).__name__)
            continue

        if extracted is None:
            reason = UNREADABLE.get(row["mime_type"] or "", f"no reader for {row['mime_type']}")
            if not dry_run:
                store.upsert_extraction(
                    db, drive_id, status="unsupported", error=reason, extracted_at=stamp
                )
            result.unsupported += 1
            _reason(result, reason)
            continue

        result.extracted += 1
        result.chars += len(extracted.text)
        result.scan_pages += extracted.scan_pages
        if extracted.method == "pymupdf" and extracted.pages:
            result.pdf_pages += extracted.pages
            result.pdf_profile.append(
                (
                    titles.get(drive_id) or drive_id,
                    len(extracted.text) / extracted.pages,
                    extracted.scan_pages,
                )
            )

        if dry_run:
            continue

        relative = f"{TEXT_SUBDIR}/{drive_id}.txt"
        destination = text_dir / f"{drive_id}.txt"
        destination.parent.mkdir(parents=True, exist_ok=True)
        destination.write_text(extracted.text, encoding="utf-8")

        store.upsert_extraction(
            db,
            drive_id,
            status="ok",
            text_path=relative,
            method=extracted.method,
            pages=extracted.pages,
            chars=len(extracted.text),
            scan_pages=extracted.scan_pages,
            ocr_pages=extracted.ocr_pages,
            extracted_at=stamp,
            error=None,
        )

    return result
