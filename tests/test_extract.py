"""Turning fetched bytes into text.

Fixtures are built here rather than committed: the five real PDFs the Phase 0
probe downloaded live in `sample/`, which is gitignored, so a test that read
them would pass on this laptop and fail everywhere else. PyMuPDF can write the
two shapes that matter -- a page whose text is text, and a page whose text is a
picture of text -- so both are synthesised.

Nothing here touches the network or Drive.
"""

from __future__ import annotations

import json
import zipfile

import pymupdf
import pytest

from agent.classroom.models import Course, Material
from agent.config import Config
from agent.db import store
from agent.files import extract

# --------------------------------------------------------------------------
# fixtures
# --------------------------------------------------------------------------


@pytest.fixture
def config(tmp_path) -> Config:
    data_dir = tmp_path / "data"
    (data_dir / "library").mkdir(parents=True)
    return Config(
        account="someone@example.com",
        timezone="Africa/Tunis",
        data_dir=data_dir,
        tracked_courses=["c1"],
        ignored_courses=[],
    )


@pytest.fixture
def conn(config):
    connection = store.connect(config.db_path)
    store.upsert_course(
        connection,
        Course(
            id="c1",
            name="Operating Systems",
            section=None,
            room=None,
            owner_id=None,
            course_state="ACTIVE",
            enrollment_code=None,
            alternate_link=None,
            creation_time=None,
            update_time=None,
            content_hash="h",
        ),
    )
    yield connection
    connection.close()


def native_pdf(path, pages=2, body="Process synchronisation. " * 20):
    """A PDF whose pages carry real text."""
    document = pymupdf.open()
    for _ in range(pages):
        page = document.new_page()
        page.insert_textbox(pymupdf.Rect(40, 40, 550, 780), body, fontsize=11)
    document.save(path)
    document.close()
    return path


def scanned_pdf(path, pages=1):
    """A PDF whose pages are an image of text, which is what a scan is.

    Rendered from a text page and re-inserted as a pixmap, so the bytes look
    exactly like the "Microsoft: Print To PDF" scans in this account: an image,
    no extractable text.
    """
    source = pymupdf.open()
    page = source.new_page()
    page.insert_textbox(pymupdf.Rect(40, 40, 550, 780), "Exercise 4. " * 40, fontsize=14)
    pixmap = page.get_pixmap(dpi=72)
    source.close()

    document = pymupdf.open()
    for _ in range(pages):
        target = document.new_page()
        target.insert_image(target.rect, pixmap=pixmap)
    document.save(path)
    document.close()
    return path


def blank_pdf(path, pages=1):
    """A section divider: no text, nothing drawn. OCR would find nothing."""
    document = pymupdf.open()
    for _ in range(pages):
        document.new_page()
    document.save(path)
    document.close()
    return path


def a_material(drive_id, title="Chapter 6.pdf") -> Material:
    return Material(
        id=f"coursework_material:p1:driveFile:{drive_id}",
        parent_type="coursework_material",
        parent_id="p1",
        course_id="c1",
        kind="driveFile",
        ref=drive_id,
        drive_id=drive_id,
        title=title,
        url=None,
        content_hash="h",
    )


def a_fetched_file(config, conn, drive_id, filename, mime_type, *, payload=None, builder=None):
    """Put a file in the library and record it as fetched, as drive.fetch would."""
    destination = config.library_dir / "files" / filename
    destination.parent.mkdir(parents=True, exist_ok=True)
    if builder is not None:
        builder(destination)
    else:
        destination.write_bytes(payload or b"")
    store.upsert_material(conn, a_material(drive_id, title=filename))
    store.upsert_extraction(
        conn,
        drive_id,
        status="fetched",
        mime_type=mime_type,
        local_path=f"files/{filename}",
    )
    return destination


# --------------------------------------------------------------------------
# decoding
# --------------------------------------------------------------------------


@pytest.mark.parametrize(
    "payload,expected",
    [
        ("plain ascii".encode("utf-8"), "plain ascii"),
        ("café naïve".encode("utf-8"), "café naïve"),
        ("with a bom".encode("utf-8-sig"), "with a bom"),
    ],
)
def test_text_decodes_across_encodings(payload, expected):
    assert extract.decode(payload) == expected


def test_a_cp1252_apostrophe_does_not_kill_the_run():
    """The 14 text/plain files have unknown encodings; utf-8 alone dies on one byte."""
    payload = "it’s a trap".encode("cp1252")
    with pytest.raises(UnicodeDecodeError):
        payload.decode("utf-8")

    assert "s a trap" in extract.decode(payload)


# --------------------------------------------------------------------------
# PDF page classification
# --------------------------------------------------------------------------


def test_a_text_page_is_native(tmp_path):
    result = extract.extract_pdf(native_pdf(tmp_path / "a.pdf", pages=3))

    assert result.pages == 3
    assert result.scan_pages == 0
    assert "synchronisation" in result.text


def test_an_image_page_is_a_scan(tmp_path):
    """Low text plus an image is the signature of every scanned sample measured."""
    result = extract.extract_pdf(scanned_pdf(tmp_path / "b.pdf", pages=2))

    assert result.pages == 2
    assert result.scan_pages == 2


def test_a_blank_page_is_not_treated_as_a_scan(tmp_path):
    """The rule is not "low text means OCR". Dividers would multiply the work."""
    result = extract.extract_pdf(blank_pdf(tmp_path / "c.pdf", pages=4))

    assert result.pages == 4
    assert result.scan_pages == 0


def test_a_mixed_pdf_is_judged_page_by_page(tmp_path):
    """Typed notes with scanned exercises stapled on is a real shape in this library."""
    document = pymupdf.open()
    page = document.new_page()
    page.insert_textbox(pymupdf.Rect(40, 40, 550, 780), "Chapter 5. " * 30, fontsize=11)

    source = pymupdf.open()
    temp = source.new_page()
    temp.insert_textbox(pymupdf.Rect(40, 40, 550, 780), "Exercise. " * 40, fontsize=14)
    pixmap = temp.get_pixmap(dpi=72)
    source.close()
    scanned = document.new_page()
    scanned.insert_image(scanned.rect, pixmap=pixmap)

    path = tmp_path / "mixed.pdf"
    document.save(path)
    document.close()

    result = extract.extract_pdf(path)

    assert result.pages == 2
    assert result.scan_pages == 1  # not 0 and not 2
    assert "Chapter 5" in result.text


def test_pages_are_separated_by_a_form_feed(tmp_path):
    result = extract.extract_pdf(native_pdf(tmp_path / "d.pdf", pages=3))

    assert result.text.count(extract.PAGE_BREAK) == 2


# --------------------------------------------------------------------------
# other formats
# --------------------------------------------------------------------------


def test_a_docx_yields_paragraphs_and_table_cells(tmp_path):
    from docx import Document

    document = Document()
    document.add_paragraph("Deadlock conditions")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Mutual exclusion"
    table.rows[0].cells[1].text = "Hold and wait"
    path = tmp_path / "notes.docx"
    document.save(path)

    result = extract.extract_docx(path)

    assert "Deadlock conditions" in result.text
    # Lab sheets put the real content in tables, invisible to `paragraphs`.
    assert "Mutual exclusion" in result.text
    assert "Hold and wait" in result.text


def test_a_pptx_yields_slide_text_and_notes(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Paging and segmentation"
    slide.notes_slide.notes_text_frame.text = "mention the TLB"
    path = tmp_path / "deck.pptx"
    presentation.save(path)

    result = extract.extract_pptx(path)

    assert "Paging and segmentation" in result.text
    assert "mention the TLB" in result.text


def test_a_notebook_yields_cells_without_outputs(tmp_path):
    notebook = {
        "cells": [
            {"cell_type": "markdown", "source": ["# Gradient descent\n"]},
            {"cell_type": "code", "source": ["import numpy as np\n"], "outputs": ["noise"]},
        ]
    }
    payload = json.dumps(notebook).encode("utf-8")

    result = extract.extract_notebook(payload)

    assert "Gradient descent" in result.text
    assert "import numpy" in result.text
    assert "noise" not in result.text


def test_json_that_is_not_a_notebook_is_not_extracted(tmp_path):
    path = tmp_path / "config.json"
    path.write_bytes(json.dumps({"setting": 1}).encode("utf-8"))

    assert extract.extract_file(path, "application/json") is None


@pytest.mark.parametrize(
    "mime_type",
    [
        "application/zip",
        "application/vnd.ms-powerpoint",  # legacy .ppt, and one is tracked
        "application/msword",
        "application/seb",
        "video/mp4",
    ],
)
def test_unreadable_formats_return_none_rather_than_raising(tmp_path, mime_type):
    path = tmp_path / "thing"
    path.write_bytes(b"not really any of these")

    assert extract.extract_file(path, mime_type) is None


def test_a_legacy_ppt_never_reaches_python_pptx(tmp_path):
    """It is an OLE file, not a zip. python-pptx raises something baffling on it."""
    path = tmp_path / "old.ppt"
    path.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1")

    assert extract.extract_file(path, "application/vnd.ms-powerpoint") is None


def test_an_unknown_mime_type_is_unsupported_not_a_crash(tmp_path):
    path = tmp_path / "mystery"
    path.write_bytes(b"\x00\x01\x02")

    assert extract.extract_file(path, "application/x-unheard-of") is None


# --------------------------------------------------------------------------
# the extract run
# --------------------------------------------------------------------------


def test_text_lands_in_the_library_and_the_row_records_it(config, conn):
    a_fetched_file(config, conn, "f1", "f1.pdf", "application/pdf", builder=native_pdf)

    result = extract.extract(config, conn)

    assert result.extracted == 1
    row = store.get_extraction(conn, "f1")
    assert row["status"] == "ok"
    assert row["text_path"] == "text/f1.txt"
    assert row["method"] == "pymupdf"
    assert row["chars"] > 0
    written = config.library_dir / "text" / "f1.txt"
    assert "synchronisation" in written.read_text(encoding="utf-8")


def test_a_second_extract_reads_nothing_again(config, conn):
    a_fetched_file(config, conn, "f1", "f1.pdf", "application/pdf", builder=native_pdf)
    extract.extract(config, conn)

    result = extract.extract(config, conn)

    assert result.extracted == 0
    assert result.skipped == 1


def test_force_re_reads_an_already_extracted_file(config, conn):
    a_fetched_file(config, conn, "f1", "f1.pdf", "application/pdf", builder=native_pdf)
    extract.extract(config, conn)

    assert extract.extract(config, conn, force=True).extracted == 1


def test_an_unsupported_file_is_recorded_with_its_reason(config, conn):
    path = a_fetched_file(config, conn, "z1", "z1.zip", "application/zip")
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("solution.py", "print(1)")

    result = extract.extract(config, conn)

    assert result.unsupported == 1
    row = store.get_extraction(conn, "z1")
    assert row["status"] == "unsupported"
    assert "zip" in row["error"]


def test_one_broken_file_does_not_end_the_run(config, conn):
    """A truncated PDF must cost one row, not the other ninety-seven files."""
    a_fetched_file(config, conn, "bad", "bad.pdf", "application/pdf", payload=b"%PDF-1.7 truncated")
    a_fetched_file(config, conn, "good", "good.pdf", "application/pdf", builder=native_pdf)

    result = extract.extract(config, conn)

    assert result.errors == 1
    assert result.extracted == 1
    assert store.get_extraction(conn, "bad")["status"] == "error"


def test_a_vanished_local_file_is_recorded_not_raised(config, conn):
    path = a_fetched_file(config, conn, "f1", "f1.pdf", "application/pdf", builder=native_pdf)
    path.unlink()

    result = extract.extract(config, conn)

    assert result.errors == 1
    assert "missing" in store.get_extraction(conn, "f1")["error"]


def test_dry_run_measures_without_writing(config, conn):
    a_fetched_file(config, conn, "f1", "f1.pdf", "application/pdf", builder=native_pdf)

    result = extract.extract(config, conn, dry_run=True)

    assert result.extracted == 1
    assert result.pdf_profile  # the measurement is still produced
    assert not (config.library_dir / "text").exists()
    assert store.get_extraction(conn, "f1")["status"] == "fetched"


def test_the_dry_run_counts_the_pages_ocr_would_have_to_read(config, conn):
    """This total is what the OCR go/no-go decision turns on."""
    a_fetched_file(config, conn, "s1", "s1.pdf", "application/pdf",
                   builder=lambda p: scanned_pdf(p, pages=3))
    a_fetched_file(config, conn, "n1", "n1.pdf", "application/pdf",
                   builder=lambda p: native_pdf(p, pages=2))

    result = extract.extract(config, conn, dry_run=True)

    assert result.scan_pages == 3
    assert result.pdf_pages == 5


def test_scan_pages_are_stored_so_the_gap_stays_visible(config, conn):
    a_fetched_file(config, conn, "s1", "s1.pdf", "application/pdf", builder=scanned_pdf)

    extract.extract(config, conn)

    row = store.get_extraction(conn, "s1")
    # Extraction succeeded and produced almost nothing. Both facts are recorded,
    # because a status of 'ok' alone would overstate what the library holds.
    assert row["status"] == "ok"
    assert row["scan_pages"] == 1
    assert row["ocr_pages"] == 0


def test_files_that_were_never_fetched_are_not_candidates(config, conn):
    store.upsert_material(conn, a_material("dead"))
    store.upsert_extraction(conn, "dead", status="trashed")

    assert extract.extract(config, conn).candidates == 0
